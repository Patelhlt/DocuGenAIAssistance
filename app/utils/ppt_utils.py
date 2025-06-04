from pptx import Presentation
import os

def parse_ppt(file_path):
    prs = Presentation(file_path)
    slides = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        if text.strip():  # Only add non-empty slides
            metadata = {
                "source": os.path.basename(file_path),
                "ext": os.path.splitext(file_path)[1].lower(),
                "type": "ppt" if file_path.lower().endswith(".ppt") else "pptx",
                "slide": slide_num
            }
            slides.append((text.strip(), metadata))

    return slides